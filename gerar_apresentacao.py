"""Gera a apresentação RACE_Apresentacao_Conclusoes.pptx — foco nas conclusões."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ──────────────────────────────────────────────────────────
AZUL_ESCURO   = RGBColor(0x1A, 0x37, 0x5E)   # título / fundo header
AZUL_MEDIO    = RGBColor(0x24, 0x5C, 0x9E)   # destaque
AZUL_CLARO    = RGBColor(0xD6, 0xE4, 0xF7)   # fundo suave
VERDE         = RGBColor(0x1A, 0x7A, 0x4A)   # positivo
VERMELHO      = RGBColor(0xC0, 0x39, 0x2B)   # negativo / limitação
CINZA_TEXTO   = RGBColor(0x33, 0x33, 0x33)
BRANCO        = RGBColor(0xFF, 0xFF, 0xFF)
AMARELO       = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout completamente em branco


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=CINZA_TEXTO,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Barra azul escura no topo com título branco."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_color=AZUL_ESCURO)
    add_text_box(slide, title, 0.3, 0.1, 12.5, 0.65,
                 font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.72, 12.5, 0.38,
                     font_size=14, bold=False, color=RGBColor(0xBF, 0xD7, 0xFF),
                     align=PP_ALIGN.LEFT)


def footer(slide, txt="RACE — Reconhecimento e Avaliação Computacional de Exercícios | UNISO 2026"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill_color=AZUL_ESCURO)
    add_text_box(slide, txt, 0.2, 7.2, 12.9, 0.3,
                 font_size=9, color=RGBColor(0xBF, 0xD7, 0xFF), align=PP_ALIGN.CENTER)


def bullet_box(slide, items, l, t, w, h,
               font_size=16, color=CINZA_TEXTO, bullet="▸ ", bold_first=False):
    """Caixa com lista de bullets (texto simples)."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        if bold_first:
            run.font.bold = True
            bold_first = False


def metric_box(slide, label, value, l, t, w=2.6, h=1.1,
               value_color=AZUL_MEDIO, bg=AZUL_CLARO):
    add_rect(slide, l, t, w, h, fill_color=bg,
             line_color=AZUL_MEDIO, line_width=Pt(1.5))
    add_text_box(slide, value, l, t + 0.05, w, 0.6,
                 font_size=30, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l, t + 0.65, w, 0.4,
                 font_size=11, bold=False, color=CINZA_TEXTO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Capa
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

add_rect(sl, 0, 0, 13.33, 7.5, fill_color=AZUL_ESCURO)
add_rect(sl, 0, 5.2, 13.33, 2.3, fill_color=RGBColor(0x12, 0x27, 0x43))

add_text_box(sl, "RACE", 0.5, 0.5, 12, 2.0,
             font_size=96, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

add_text_box(sl, "Reconhecimento e Avaliação Computacional de Exercícios",
             0.5, 2.5, 12, 0.65,
             font_size=22, color=RGBColor(0xBF, 0xD7, 0xFF), align=PP_ALIGN.CENTER)

add_text_box(sl, "Conclusões", 0.5, 3.2, 12, 0.65,
             font_size=30, bold=True, color=AMARELO, align=PP_ALIGN.CENTER)

add_text_box(sl,
             "Rafael Luckner Flora  ·  Leticia Ruivo Tambelli  ·  Paulo Vitor Nascimento Silva  ·  Rodrigo Araújo Sousa",
             0.5, 5.4, 12, 0.45,
             font_size=13, color=RGBColor(0x90, 0xB4, 0xE0), align=PP_ALIGN.CENTER)

add_text_box(sl, "Orientador: Prof. Dr. Jaime Ranulfo Leite Filho",
             0.5, 5.88, 12, 0.35,
             font_size=12, color=RGBColor(0x90, 0xB4, 0xE0), align=PP_ALIGN.CENTER)

add_text_box(sl, "Visão Computacional e Reconhecimento de Imagens  ·  UNISO  ·  2026",
             0.5, 6.85, 12, 0.35,
             font_size=11, color=RGBColor(0x70, 0x96, 0xC0), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — O problema e a proposta (contexto em 1 slide compacto)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "O que o RACE se propôs a resolver?")
footer(sl)

# Bloco problema
add_rect(sl, 0.3, 1.3, 5.9, 5.7, fill_color=RGBColor(0xFF, 0xF5, 0xEE),
         line_color=AMARELO, line_width=Pt(2))
add_text_box(sl, "Problema", 0.55, 1.42, 5.4, 0.45,
             font_size=18, bold=True, color=RGBColor(0x9A, 0x52, 0x00))
bullet_box(sl, [
    "Monitorar exercícios exige\n   fisioterapeuta ou personal trainer",
    "Sensores dedicados são caros\n   e de difícil acesso",
    "Soluções existentes exigem\n   marcadores corporais ou câmeras especiais",
], 0.55, 1.92, 5.55, 3.5, font_size=14, color=CINZA_TEXTO)

# Seta
add_text_box(sl, "→", 6.3, 3.6, 0.7, 0.7,
             font_size=36, bold=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

# Bloco proposta
add_rect(sl, 7.1, 1.3, 5.9, 5.7, fill_color=RGBColor(0xF0, 0xF8, 0xFF),
         line_color=AZUL_MEDIO, line_width=Pt(2))
add_text_box(sl, "Proposta: RACE", 7.35, 1.42, 5.4, 0.45,
             font_size=18, bold=True, color=AZUL_ESCURO)
bullet_box(sl, [
    "Vídeo comum como única entrada",
    "MediaPipe Pose: 33 landmarks,\n   sem marcadores, em tempo real",
    "8 ângulos articulares bilaterais\n   ponderados por visibilidade",
    "Random Forest (janelas 15 frames)\n   + máquina de estados para contagem",
    "Sem sensores  ·  Sem calibração",
], 7.35, 1.92, 5.55, 4.5, font_size=14, color=CINZA_TEXTO)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — Conclusão 1: Perfis angulares
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Conclusão 1 — A abordagem funciona",
           "Ângulos articulares como representação do movimento humano")
footer(sl)

add_rect(sl, 0.3, 1.3, 12.73, 0.7, fill_color=RGBColor(0xD4, 0xED, 0xDA),
         line_color=VERDE, line_width=Pt(1))
add_text_box(sl,
    "Os perfis angulares extraídos pelo MediaPipe apresentam assinaturas biomecânicas "
    "distintas e reprodutíveis para cada exercício, mesmo com variações de câmera, "
    "iluminação e sujeito.",
    0.5, 1.33, 12.3, 0.65, font_size=13, bold=True, color=RGBColor(0x1A, 0x5C, 0x34))

exerc2 = [
    ("Agachamento",
     "Joelho: 50°–160°  |  Quadril: 25°–175°\nAmplitude ~110° (joelho) e ~150° (quadril)\nAlta simetria bilateral  ·  Ciclos regulares por 177 s"),
    ("Flexão de braço",
     "Cotovelo: distribuição bimodal\n(posição estendida ↔ contraída)\nJoelho/Quadril estáveis ~175° → confirma posição deitada"),
    ("Rosca bíceps",
     "Cotovelo: 80°–170°\nAssimetria bilateral consistente (dominância)\nFadiga visível após ~115 s: amplitude cai, irregularidade sobe"),
]

for i, (titulo, desc) in enumerate(exerc2):
    lx = 0.3 + i * 4.35
    add_rect(sl, lx, 2.15, 4.1, 4.6, fill_color=RGBColor(0xF0, 0xF7, 0xFF),
             line_color=AZUL_MEDIO, line_width=Pt(1.5))
    add_rect(sl, lx, 2.15, 4.1, 0.5, fill_color=AZUL_ESCURO)
    add_text_box(sl, titulo, lx + 0.1, 2.17, 3.9, 0.45,
                 font_size=15, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text_box(sl, desc, lx + 0.2, 2.75, 3.75, 3.8,
                 font_size=13, color=CINZA_TEXTO)

add_text_box(sl,
    "O sinal bilateral ponderado por visibilidade reduz o impacto de oclusões e variações de orientação da câmera.",
    0.3, 6.9, 12.73, 0.38,
    font_size=12, italic=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Conclusão 2: Classificação
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Conclusão 2 — O modelo generaliza para novos sujeitos",
           "Random Forest treinado no Sujeito 1, avaliado nos Sujeitos 2 e 3")
footer(sl)

add_rect(sl, 0.3, 1.3, 12.73, 0.7, fill_color=RGBColor(0xD4, 0xED, 0xDA),
         line_color=VERDE, line_width=Pt(1))
add_text_box(sl,
    "O RF atingiu 95,7% de recall médio nos exercícios ativos em sujeitos nunca vistos "
    "durante o treinamento — superando o baseline heurístico em 5,4 pontos percentuais.",
    0.5, 1.33, 12.3, 0.65, font_size=13, bold=True, color=RGBColor(0x1A, 0x5C, 0x34))

# Métricas destaque
metric_box(sl, "Recall médio\nexercícios ativos — BL", "90,3 %",
           0.4, 2.15, w=3.1, h=1.25, value_color=AZUL_MEDIO)
add_text_box(sl, "+5,4 pp →", 3.55, 2.55, 1.4, 0.5,
             font_size=20, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
metric_box(sl, "Recall médio\nexercícios ativos — RF", "95,7 %",
           5.0, 2.15, w=3.1, h=1.25, value_color=VERDE, bg=RGBColor(0xD4, 0xED, 0xDA))

# Tabela por classe
headers = ["Exercício", "BL Recall", "RF Recall", "BL Precisão", "RF Precisão"]
rows = [
    ("Flexão",         "0,871", "0,920", "0,584", "0,700"),
    ("Agachamento",    "0,982", "0,950", "0,809", "1,000"),
    ("Descanso ⚠",    "0,185", "0,110", "1,000", "0,460"),
    ("Rosca bíceps",   "0,863", "1,000", "0,777", "0,670"),
    ("Média c/ desc.", "0,725", "0,745", "0,793", "0,708"),
    ("Média ativos ★", "0,905", "0,957", "0,724", "0,790"),
]
col_w = [2.6, 1.7, 1.7, 1.8, 1.8]
starts = [0.4]
for w in col_w[:-1]:
    starts.append(starts[-1] + w)

for ci, (hd, sx, cw) in enumerate(zip(headers, starts, col_w)):
    add_rect(sl, sx, 3.55, cw, 0.38, fill_color=AZUL_ESCURO)
    add_text_box(sl, hd, sx + 0.05, 3.57, cw - 0.1, 0.34,
                 font_size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    ty = 3.95 + ri * 0.4
    is_media = ri >= 4
    is_descanso = "descanso" in row[0].lower()
    bg = RGBColor(0xD4, 0xED, 0xDA) if ri == 5 else (
         RGBColor(0xFF, 0xF0, 0xF0) if is_descanso else (
         AZUL_CLARO if ri % 2 == 0 else BRANCO))
    for ci, (cell, sx, cw) in enumerate(zip(row, starts, col_w)):
        add_rect(sl, sx, ty, cw, 0.4, fill_color=bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        cc = CINZA_TEXTO
        if ci == 2:
            try:
                if float(row[2].replace(',','.')) > float(row[1].replace(',','.')):
                    cc = VERDE
                elif float(row[2].replace(',','.')) < float(row[1].replace(',','.')):
                    cc = VERMELHO
            except: pass
        add_text_box(sl, cell, sx + 0.05, ty + 0.05, cw - 0.1, 0.32,
                     font_size=11, bold=(ci == 0 or ri >= 4), color=cc,
                     align=PP_ALIGN.CENTER)

add_text_box(sl,
    "Descanso = principal gargalo: movimentos livres de braço sobrepõem as faixas angulares dos exercícios ativos",
    0.4, 7.1, 12.5, 0.3, font_size=11, italic=True, color=VERMELHO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Conclusão 3: Contagem de repetições
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Conclusão 3 — A contagem automática é robusta",
           "Máquina de estados com histerese — sem treinamento prévio")
footer(sl)

add_rect(sl, 0.3, 1.3, 12.73, 0.7, fill_color=RGBColor(0xD4, 0xED, 0xDA),
         line_color=VERDE, line_width=Pt(1))
add_text_box(sl,
    "A máquina de estados detectou ciclos completos de movimento (ESTENDIDO → CONTRAÍDO → ESTENDIDO) "
    "de forma coerente com a cadência real dos vídeos, sem qualquer treinamento.",
    0.5, 1.33, 12.3, 0.65, font_size=13, bold=True, color=RGBColor(0x1A, 0x5C, 0x34))

# Como funciona
add_rect(sl, 0.3, 2.15, 5.9, 2.7, fill_color=RGBColor(0xF0, 0xF5, 0xFF),
         line_color=AZUL_MEDIO, line_width=Pt(1.5))
add_text_box(sl, "Como funciona", 0.5, 2.22, 5.5, 0.45,
             font_size=16, bold=True, color=AZUL_ESCURO)
bullet_box(sl, [
    "Sinal bilateral ponderado por visibilidade",
    "2 estados: ESTENDIDO e CONTRAÍDO",
    "Histerese: limiares distintos de\n   entrada e saída → sem dupla contagem",
    "Repetição = ciclo completo concluído",
], 0.5, 2.7, 5.6, 1.95, font_size=13)

# Tabela de resultados
add_rect(sl, 6.5, 2.15, 6.5, 2.7, fill_color=RGBColor(0xF0, 0xFB, 0xF0),
         line_color=VERDE, line_width=Pt(1.5))
add_text_box(sl, "Resultados — Sujeito 1", 6.7, 2.22, 6.1, 0.45,
             font_size=16, bold=True, color=RGBColor(0x1A, 0x5C, 0x34))

t_headers = ["Exercício", "Reps", "Reps/min", "Duração"]
t_rows = [
    ("Flexão de braço", "37", "38,9 /min", "57 s"),
    ("Agachamento",     "54", "18,3 /min", "177 s"),
    ("Rosca bíceps",    "12", " 6,0 /min", "121 s"),
]
tcol_w = [2.5, 0.9, 1.5, 1.2]
tstarts = [6.6]
for w in tcol_w[:-1]:
    tstarts.append(tstarts[-1] + w)

for ci, (hd, sx, cw) in enumerate(zip(t_headers, tstarts, tcol_w)):
    add_rect(sl, sx, 2.72, cw, 0.36, fill_color=AZUL_ESCURO)
    add_text_box(sl, hd, sx+0.04, 2.74, cw-0.08, 0.3,
                 font_size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(t_rows):
    ty = 3.1 + ri * 0.42
    bg = RGBColor(0xD4, 0xED, 0xDA) if ri % 2 == 0 else BRANCO
    for ci, (cell, sx, cw) in enumerate(zip(row, tstarts, tcol_w)):
        add_rect(sl, sx, ty, cw, 0.42, fill_color=bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text_box(sl, cell, sx+0.04, ty+0.06, cw-0.08, 0.32,
                     font_size=12, bold=(ci == 0), color=CINZA_TEXTO,
                     align=PP_ALIGN.CENTER)

bullet_box(sl, [
    "Contagens coerentes com as cadências observadas nos vídeos",
    "Histerese eliminou duplas contagens por ruído angular",
    "Funciona sem dados de treino — generalizável a novos exercícios",
], 0.4, 5.05, 12.5, 1.6, font_size=14, color=VERDE, bullet="✓  ")


# ════════════════════════════════════════════════════════════════════════════
# Slide 7 — Contribuição metodológica
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "O que torna o RACE diferente?",
           "Três escolhas metodológicas que sustentam os resultados")
footer(sl)

contrib = [
    (VERDE,      "Sinal bilateral ponderado por visibilidade",
     "Em vez de escolher um lado do corpo arbitrariamente, o sistema "
     "combina os dois lados dando mais peso ao que o MediaPipe vê melhor. "
     "Isso reduz o impacto de oclusões e variações de orientação da câmera."),
    (AZUL_MEDIO, "Representação temporal por janelas deslizantes",
     "Cada amostra de classificação representa 15 frames (~3 s) do movimento, "
     "não uma posição instantânea. Isso captura a dinâmica do exercício e "
     "torna o modelo robusto a ruído de frame isolado."),
    (AMARELO,    "Avaliação cruzada entre sujeitos",
     "Treino e teste usam sujeitos distintos — garantindo que o modelo aprende "
     "padrões de movimento, não características de um indivíduo específico. "
     "Isso fornece uma estimativa realista da generalização."),
]

for idx, (cor, titulo, desc) in enumerate(contrib):
    ly = 1.4 + idx * 1.75
    add_rect(sl, 0.3, ly, 0.12, 1.45, fill_color=cor)
    add_rect(sl, 0.52, ly, 12.5, 1.45, fill_color=RGBColor(0xF4, 0xF7, 0xFF),
             line_color=RGBColor(0xCC, 0xD8, 0xF0), line_width=Pt(1))
    add_text_box(sl, titulo, 0.7, ly + 0.1, 12.1, 0.45,
                 font_size=17, bold=True, color=AZUL_ESCURO)
    add_text_box(sl, desc, 0.7, ly + 0.58, 12.1, 0.8,
                 font_size=13, color=CINZA_TEXTO)


# ════════════════════════════════════════════════════════════════════════════
# Slide 8 — Limitações e Trabalhos Futuros
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Limitações & Próximos Passos")
footer(sl)

add_rect(sl, 0.3, 1.3, 6.0, 5.7, fill_color=RGBColor(0xFF, 0xF5, 0xF5),
         line_color=VERMELHO, line_width=Pt(1.5))
add_text_box(sl, "⚠  Limitações observadas", 0.5, 1.42, 5.6, 0.45,
             font_size=16, bold=True, color=VERMELHO)
bullet_box(sl, [
    "Descanso: recall muito baixo\n   (0,110 RF) — principal gargalo",
    "Apenas 3 sujeitos válidos\n   (adultos jovens, baixa diversidade)",
    "Ângulos 2D: distorção quando o\n   plano de movimento ≠ plano da imagem",
    "Apenas 3 exercícios cobertos;\n   sem equipamentos ou rotação de tronco",
    "Contagem avaliada qualitativamente\n   (sem erro absoluto calculado)",
], 0.5, 1.95, 5.65, 4.8, font_size=13, color=CINZA_TEXTO)

add_rect(sl, 6.85, 1.3, 6.1, 5.7, fill_color=RGBColor(0xF0, 0xFB, 0xF0),
         line_color=VERDE, line_width=Pt(1.5))
add_text_box(sl, "→  Próximos passos", 7.05, 1.42, 5.8, 0.45,
             font_size=16, bold=True, color=VERDE)
bullet_box(sl, [
    "Features de variância temporal\n   (desvio padrão na janela) → melhorar descanso",
    "Ampliar dataset: diversidade de\n   sujeitos, biotipos, idades e sexos",
    "Novos exercícios com equipamentos\n   e rotações de tronco",
    "Pose 3D (câmera estéreo / depth)\n   → eliminar distorção de perspectiva",
    "Feedback de qualidade: assimetria,\n   amplitude insuficiente, fadiga",
    "App móvel / web em tempo real",
], 7.05, 1.95, 5.85, 4.8, font_size=13, color=CINZA_TEXTO)


# ════════════════════════════════════════════════════════════════════════════
# Slide 9 — Obrigado
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=AZUL_ESCURO)
add_rect(sl, 0, 5.3, 13.33, 2.2, fill_color=RGBColor(0x12, 0x27, 0x43))

add_text_box(sl, "Obrigado!", 0.5, 1.3, 12.3, 1.8,
             font_size=72, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

# Síntese das 3 conclusões na capa final
for idx, (txt, cor) in enumerate([
    ("✓  Ângulos articulares capturam padrões biomecânicos distintos", VERDE),
    ("✓  RF: 95,7% recall (ativos) em sujeitos novos",                VERDE),
    ("✓  Contagem automática coerente sem treinamento",               VERDE),
]):
    add_text_box(sl, txt, 1.5, 3.15 + idx * 0.52, 10.5, 0.48,
                 font_size=16, bold=True, color=cor, align=PP_ALIGN.CENTER)

add_text_box(sl, "Rafael Luckner Flora  ·  Leticia Ruivo Tambelli\nPaulo Vitor Nascimento Silva  ·  Rodrigo Araújo Sousa",
             0.5, 5.5, 12.3, 0.7,
             font_size=13, color=RGBColor(0x90, 0xB4, 0xE0), align=PP_ALIGN.CENTER)

add_text_box(sl, "Orientador: Prof. Dr. Jaime Ranulfo Leite Filho  ·  UNISO 2026",
             0.5, 6.6, 12.3, 0.35,
             font_size=12, color=RGBColor(0x70, 0x96, 0xC0), align=PP_ALIGN.CENTER)


# ── Salvar ───────────────────────────────────────────────────────────────────
out = "RACE_Apresentacao_Conclusoes.pptx"
prs.save(out)
print(f"✓ Apresentação salva: {out}  ({prs.slides.__len__()} slides)")
import sys; sys.exit(0)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — Dados e Metodologia (orphaned)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Dados e Metodologia", "Divisão por sujeito — sem vazamento de dados")
footer(sl)

# Coluna esquerda — Dataset
add_rect(sl, 0.3, 1.3, 6.0, 5.7, fill_color=RGBColor(0xF5, 0xF8, 0xFF),
         line_color=AZUL_CLARO, line_width=Pt(1))
add_text_box(sl, "Dataset", 0.5, 1.45, 5.7, 0.45,
             font_size=17, bold=True, color=AZUL_ESCURO)
bullet_box(sl, [
    "4 exercícios: flexão, agachamento,\n   rosca bíceps, descanso",
    "3 sujeitos com dados válidos",
    "Gravações a 5 fps (modelo full)",
    "Treino: Sujeito 1  |  Teste: Sujeitos 2 e 3",
    "Balanceamento: undersampling por split",
], 0.5, 1.95, 5.7, 4.4, font_size=14)

# Coluna direita — Features
add_rect(sl, 6.8, 1.3, 6.2, 5.7, fill_color=RGBColor(0xF5, 0xF8, 0xFF),
         line_color=AZUL_CLARO, line_width=Pt(1))
add_text_box(sl, "Engenharia de features", 7.0, 1.45, 5.9, 0.45,
             font_size=17, bold=True, color=AZUL_ESCURO)
bullet_box(sl, [
    "12 landmarks usados:\n   ombros, cotovelos, pulsos,\n   quadris, joelhos e tornozelos",
    "8 ângulos / frame (4 pares bilaterais):\n   cotovelo, ombro, joelho, quadril",
    "Janela deslizante: 15 frames (~3 s)",
    "Stride = 1  →  120 features / janela",
    "StandardScaler (reajustado só no treino)",
], 7.0, 1.95, 5.9, 4.4, font_size=14)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — Experimento 1: Perfis Angulares
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Experimento 1 — Perfis Angulares Articulares",
           "Validação biomecânica dos sinais extraídos")
footer(sl)

exerc = [
    ("Agachamento",
     ["Joelho: 50°–160° (amplitude ~110°)",
      "Quadril: 25°–175° (amplitude ~150°)",
      "Alta simetria bilateral direito/esquerdo",
      "Padrão periódico mantido por ~180 s"]),
    ("Flexão de braço",
     ["Cotovelo: distribuição bimodal\n   (posição estendida ↔ contraída)",
      "Joelho/Quadril: estáveis (~175°)\n   confirma posição deitada",
      "Ciclos regulares, boa simetria"]),
    ("Rosca bíceps",
     ["Cotovelo principal: 80°–170°",
      "Assimetria bilateral consistente\n   (dominância lateral)",
      "Sinal de fadiga após ~115 s:\n   amplitude reduz, irregularidade aumenta"]),
]
cols = [0.3, 4.55, 8.8]
for i, (titulo, pontos) in enumerate(exerc):
    lx = cols[i]
    add_rect(sl, lx, 1.3, 4.1, 5.7, fill_color=RGBColor(0xF0, 0xF5, 0xFF),
             line_color=AZUL_MEDIO, line_width=Pt(1.5))
    add_text_box(sl, titulo, lx + 0.1, 1.4, 3.9, 0.5,
                 font_size=16, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    bullet_box(sl, pontos, lx + 0.15, 1.95, 3.8, 4.8, font_size=13)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Experimento 2: Resultados de Classificação
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Experimento 2 — Classificação de Exercícios",
           "Baseline heurístico vs. Random Forest (conjunto de teste)")
footer(sl)

# Métricas globais
add_text_box(sl, "Recall médio — exercícios ativos (excl. descanso)",
             0.4, 1.3, 12.5, 0.4, font_size=14, bold=True, color=AZUL_ESCURO)

metric_box(sl, "Baseline heurístico", "90,3 %", 0.5, 1.75, w=3.0, h=1.1,
           value_color=AZUL_MEDIO)
add_text_box(sl, "→  +5,4 pp", 3.55, 2.05, 1.5, 0.5,
             font_size=18, bold=True, color=VERDE, align=PP_ALIGN.CENTER)
metric_box(sl, "Random Forest", "95,7 %", 5.1, 1.75, w=3.0, h=1.1,
           value_color=VERDE, bg=RGBColor(0xD4, 0xED, 0xDA))

metric_box(sl, "Recall médio macro (c/ descanso) — BL", "72,5 %",
           0.5, 3.0, w=2.8, h=0.95)
metric_box(sl, "Recall médio macro (c/ descanso) — RF", "74,5 %",
           3.4, 3.0, w=2.8, h=0.95)

# Tabela por classe
headers = ["Exercício", "BL Recall", "RF Recall", "BL Precisão", "RF Precisão"]
rows = [
    ("Flexão",        "0,871", "0,920", "0,584", "0,700"),
    ("Agachamento",   "0,982", "0,950", "0,809", "1,000"),
    ("Descanso",      "0,185", "0,110", "1,000", "0,460"),
    ("Rosca bíceps",  "0,863", "1,000", "0,777", "0,670"),
    ("Média (c/ desc.)", "0,725", "0,745", "0,793", "0,708"),
    ("Média (ativos)", "0,905", "0,957", "0,724", "0,790"),
]

col_w = [2.2, 1.5, 1.5, 1.6, 1.6]
col_x = [8.7, 10.9, 12.4, 13.9, 15.5]  # ajustados para caber
# reposicionar para largura da tela
total_w = 12.0
col_w = [2.4, 1.7, 1.7, 1.8, 1.8]
starts = [0.5]
for w in col_w[:-1]:
    starts.append(starts[-1] + w)

# header da tabela
for ci, (hd, sx, cw) in enumerate(zip(headers, starts, col_w)):
    add_rect(sl, sx, 3.1, cw, 0.38, fill_color=AZUL_ESCURO)
    add_text_box(sl, hd, sx + 0.05, 3.12, cw - 0.1, 0.35,
                 font_size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

row_h = 0.38
for ri, row in enumerate(rows):
    ty = 3.5 + ri * row_h
    is_last = ri >= 4
    bg = RGBColor(0xD4, 0xED, 0xDA) if is_last else (
         AZUL_CLARO if ri % 2 == 0 else BRANCO)
    for ci, (cell, sx, cw) in enumerate(zip(row, starts, col_w)):
        add_rect(sl, sx, ty, cw, row_h, fill_color=bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        # destaque verde para RF melhor que BL em recall
        cell_color = CINZA_TEXTO
        if ci == 2 and float(row[2].replace(',', '.')) >= float(row[1].replace(',', '.')):
            cell_color = VERDE
        if ci == 2 and row[0] == "Descanso":
            cell_color = VERMELHO
        bold_cell = is_last or ci == 0
        add_text_box(sl, cell, sx + 0.05, ty + 0.04, cw - 0.1, row_h - 0.05,
                     font_size=11, bold=bold_cell, color=cell_color,
                     align=PP_ALIGN.CENTER)

add_text_box(sl, "⚠  Descanso = principal gargalo: movimentos livres sobrepõem faixas angulares dos exercícios ativos",
             0.5, 6.8, 12.3, 0.35,
             font_size=11, italic=True, color=VERMELHO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Experimento 3: Contagem de Repetições
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Experimento 3 — Contagem de Repetições",
           "Máquina de estados com histerese — sem treinamento prévio")
footer(sl)

add_text_box(sl, "Mecanismo", 0.4, 1.35, 12.5, 0.4,
             font_size=16, bold=True, color=AZUL_ESCURO)
bullet_box(sl, [
    "Sinal bilateral ponderado por visibilidade → 2 estados: ESTENDIDO e CONTRAÍDO",
    "Histerese: transição só ocorre ao cruzar limiares distintos de entrada e saída → evita dupla contagem por ruído",
    "Repetição registrada a cada ciclo ESTENDIDO → CONTRAÍDO → ESTENDIDO completo",
], 0.5, 1.78, 12.3, 1.5, font_size=13)

# Tabela de resultados
add_text_box(sl, "Resultados — Sujeito 1", 0.4, 3.4, 12.5, 0.4,
             font_size=16, bold=True, color=AZUL_ESCURO)

t_headers = ["Exercício", "Reps detectadas", "Reps/min", "Duração (s)"]
t_rows = [
    ("Flexão de braço", "37", "38,9", "57"),
    ("Agachamento",     "54", "18,3", "177"),
    ("Rosca bíceps",    "12", " 6,0", "121"),
]
tcol_w = [3.8, 2.8, 2.8, 2.8]
tstarts = [0.5]
for w in tcol_w[:-1]:
    tstarts.append(tstarts[-1] + w)

for ci, (hd, sx, cw) in enumerate(zip(t_headers, tstarts, tcol_w)):
    add_rect(sl, sx, 3.85, cw, 0.42, fill_color=AZUL_ESCURO)
    add_text_box(sl, hd, sx + 0.05, 3.87, cw - 0.1, 0.38,
                 font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(t_rows):
    ty = 4.3 + ri * 0.48
    bg = AZUL_CLARO if ri % 2 == 0 else BRANCO
    for ci, (cell, sx, cw) in enumerate(zip(row, tstarts, tcol_w)):
        add_rect(sl, sx, ty, cw, 0.48, fill_color=bg,
                 line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
        add_text_box(sl, cell, sx + 0.05, ty + 0.06, cw - 0.1, 0.38,
                     font_size=14, bold=(ci == 0), color=CINZA_TEXTO,
                     align=PP_ALIGN.CENTER)

bullet_box(sl, [
    "Contagens coerentes com cadências observadas nos vídeos",
    "Sem contagens duplicadas causadas por ruído no sinal angular",
], 0.5, 5.85, 12.3, 1.1, font_size=13, color=VERDE, bullet="✓ ")


# ════════════════════════════════════════════════════════════════════════════
# Slide 7 — Conclusões
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Conclusões", "O sistema RACE demonstrou viabilidade em três frentes")
footer(sl)

conclusoes = [
    ("Perfis angulares consistentes",
     "Padrões biomecânicos distintos por exercício, mesmo com variações de câmera e iluminação."),
    ("Classificação com generalização",
     "RF atingiu 95,7% de recall médio (exercícios ativos) em sujeitos nunca vistos no treino."),
    ("Contagem robusta de repetições",
     "Máquina de estados detectou ciclos sem treinamento, com contagens coerentes aos vídeos."),
]

for i, (titulo, desc) in enumerate(conclusoes):
    ly = 1.4 + i * 1.65
    add_rect(sl, 0.4, ly, 0.55, 1.15, fill_color=AZUL_ESCURO)
    add_text_box(sl, str(i + 1), 0.4, ly + 0.2, 0.55, 0.7,
                 font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_rect(sl, 1.05, ly, 11.85, 1.15, fill_color=RGBColor(0xF0, 0xF5, 0xFF),
             line_color=AZUL_CLARO, line_width=Pt(1))
    add_text_box(sl, titulo, 1.2, ly + 0.06, 11.5, 0.45,
                 font_size=16, bold=True, color=AZUL_ESCURO)
    add_text_box(sl, desc, 1.2, ly + 0.52, 11.5, 0.58,
                 font_size=13, color=CINZA_TEXTO)

add_text_box(sl, "Contribuição metodológica: ângulos bilaterais ponderados por visibilidade + representação temporal por janelas + avaliação cruzada entre sujeitos",
             0.4, 6.35, 12.5, 0.55,
             font_size=12, bold=True, color=AZUL_MEDIO, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# Slide 8 — Limitações e Trabalhos Futuros
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
header_bar(sl, "Limitações & Trabalhos Futuros")
footer(sl)

add_rect(sl, 0.3, 1.3, 6.15, 5.7, fill_color=RGBColor(0xFF, 0xF5, 0xF5),
         line_color=VERMELHO, line_width=Pt(1.5))
add_text_box(sl, "⚠  Limitações", 0.5, 1.4, 5.7, 0.45,
             font_size=16, bold=True, color=VERMELHO)
bullet_box(sl, [
    "Descanso: recall baixo (0,110 RF)\n   movimentos livres sobrepõem exercícios ativos",
    "Apenas 3 sujeitos — adultos jovens\n   masc.; pouca diversidade de biotipo",
    "Ângulos 2D: distorção de perspectiva\n   quando plano de mov. ≠ plano da imagem",
    "3 exercícios apenas;\n   sem equipamentos ou rotação de tronco",
    "Contagem sem validação quantitativa\n   (sem erro absoluto calculado)",
], 0.5, 1.9, 5.8, 4.8, font_size=13, color=CINZA_TEXTO)

add_rect(sl, 6.85, 1.3, 6.15, 5.7, fill_color=RGBColor(0xF0, 0xFB, 0xF0),
         line_color=VERDE, line_width=Pt(1.5))
add_text_box(sl, "→  Trabalhos futuros", 7.05, 1.4, 5.9, 0.45,
             font_size=16, bold=True, color=VERDE)
bullet_box(sl, [
    "Features de variância temporal\n   (std dentro da janela) para detectar descanso",
    "Expansão do dataset — diversidade\n   de sujeitos, biotipos e idades",
    "Ampliar repertório: novos exercícios\n   com equipamentos e rotações",
    "Estimativa de pose 3D (câmera estéreo\n   ou depth model) → eliminar distorção 2D",
    "Feedback de qualidade de execução\n   (assimetria, amplitude insuficiente, fadiga)",
    "App móvel / web em tempo real",
], 7.05, 1.9, 5.9, 4.8, font_size=13, color=CINZA_TEXTO)


# ════════════════════════════════════════════════════════════════════════════
# Slide 9 — Obrigado
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill_color=AZUL_ESCURO)
add_rect(sl, 0, 5.5, 13.33, 2.0, fill_color=RGBColor(0x12, 0x27, 0x43))

add_text_box(sl, "Obrigado!", 0.5, 1.5, 12.3, 1.8,
             font_size=72, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

add_text_box(sl, "RACE — Reconhecimento e Avaliação Computacional de Exercícios",
             0.5, 3.4, 12.3, 0.6,
             font_size=20, color=RGBColor(0xBF, 0xD7, 0xFF), align=PP_ALIGN.CENTER)

add_text_box(sl, "Rafael Luckner Flora  ·  Leticia Ruivo Tambelli\nPaulo Vitor Nascimento Silva  ·  Rodrigo Araújo Sousa",
             0.5, 4.1, 12.3, 0.8,
             font_size=14, color=RGBColor(0x90, 0xB4, 0xE0), align=PP_ALIGN.CENTER)

add_text_box(sl, "Orientador: Prof. Dr. Jaime Ranulfo Leite Filho  ·  UNISO 2026",
             0.5, 5.7, 12.3, 0.4,
             font_size=12, color=RGBColor(0x70, 0x96, 0xC0), align=PP_ALIGN.CENTER)


# ── Salvar ───────────────────────────────────────────────────────────────────
out = "RACE_Apresentacao_Conclusoes.pptx"
prs.save(out)
print(f"✓ Apresentação salva: {out}  ({prs.slides.__len__()} slides)")
